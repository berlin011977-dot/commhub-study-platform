#!/usr/bin/env python3
"""update_site.py

Runs the full content pipeline:
1️⃣ Recursively discover lecture files (.pptx, .pdf, .docx) in the project folder.
2️⃣ Extract raw slide/text content while stripping headers, footers, lecturer names, duplicate titles, "Thank You" slides, and slide numbers.
3️⃣ Clean grammar/spelling using LanguageTool.
4️⃣ Re‑structure the cleaned text into logical sections (concepts, definitions, examples, exam notes).
5️⃣ Translate every English section to Arabic (RTL) using googletrans (fallback to a local model if needed).
6️⃣ Emit a unified bilingual JSON file (`content.json`).
7️⃣ Call `generate_site.py` to rebuild the static HTML.

The script is intended to be run manually (`python update_site.py`) whenever new lecture files are added.
"""

import os, sys, json, re, argparse, subprocess
from pathlib import Path

# ---------- Helper Imports (installed on‑demand ----------
def ensure_package(pkg_name: str, import_name: str = None):
    import_name = import_name or pkg_name
    try:
        __import__(import_name)
    except ImportError:
        print(f"Installing missing package: {pkg_name} ...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", pkg_name, "-q"])
        print("Done.")

ensure_package("python-pptx", "pptx")
ensure_package("PyMuPDF", "fitz")
ensure_package("python-docx", "docx")
ensure_package("language-tool-python", "language_tool_python")
ensure_package("googletrans==4.0.0rc1", "googletrans")
# ------------------------------------------------------

import pptx
import fitz  # PyMuPDF
import docx
import language_tool_python
from googletrans import Translator

# Global regex patterns for cleaning
HEADER_FOOTER_PATTERN = re.compile(r"^(\s*\d+\s*)?$", re.IGNORECASE)
THANKYOU_PATTERN = re.compile(r"thank\s*you", re.IGNORECASE)
LECTURER_NAME_PATTERN = re.compile(r"(Prof|Dr)\.?.+", re.IGNORECASE)


def is_valid_slide(text: str) -> bool:
    """Return True if the slide contains educational content.
    Filters out empty slides, thank‑you slides, headers/footers, and lecturer names.
    """
    stripped = text.strip()
    if not stripped:
        return False
    if THANKYOU_PATTERN.search(stripped):
        return False
    if HEADER_FOOTER_PATTERN.fullmatch(stripped):
        return False
    if LECTURER_NAME_PATTERN.search(stripped):
        return False
    return True


def clean_text(text: str) -> str:
    """Apply grammar/spelling correction and strip unwanted whitespace."""
    tool = language_tool_python.LanguageTool('en-US')
    matches = tool.check(text)
    corrected = language_tool_python.utils.correct(text, matches)
    corrected = re.sub(r"\n{3,}", "\n\n", corrected)
    return corrected.strip()


def extract_from_pptx(path: Path) -> list[str]:
    prs = pptx.Presentation(str(path))
    slides_text = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text.strip()
                if txt:
                    texts.append(txt)
        slide_text = " ".join(texts)
        if is_valid_slide(slide_text):
            slides_text.append(slide_text)
    return slides_text


def extract_from_pdf(path: Path) -> list[str]:
    doc = fitz.open(str(path))
    pages_text = []
    for page in doc:
        txt = page.get_text("text").strip()
        if is_valid_slide(txt):
            pages_text.append(txt)
    return pages_text


def extract_from_docx(path: Path) -> list[str]:
    document = docx.Document(str(path))
    paragraphs = []
    for para in document.paragraphs:
        txt = para.text.strip()
        if is_valid_slide(txt):
            paragraphs.append(txt)
    return paragraphs


def discover_files(root: Path) -> list[Path]:
    patterns = ("*.pptx", "*.pdf", "*.docx")
    files = []
    for pattern in patterns:
        files.extend(root.rglob(pattern))
    return files


def group_into_sections(texts: list[str]) -> list[dict]:
    full = "\n\n".join(texts)
    raw_sections = [sec.strip() for sec in re.split(r"\n{2,}", full) if sec.strip()]
    sections = []
    for raw in raw_sections:
        lines = raw.split("\n")
        definition = ""
        examples = []
        exam_notes = []
        for line in lines:
            low = line.lower()
            if "definition" in low or "means" in low:
                definition = line.strip()
            elif "example" in low or "e.g." in low:
                examples.append(line.strip())
            elif "exam" in low or "note" in low:
                exam_notes.append(line.strip())
        sections.append({
            "raw": raw,
            "definition": definition,
            "examples": examples,
            "exam_notes": exam_notes,
        })
    return sections


def translate_section(section: dict, translator: Translator) -> dict:
    ar_section = {}
    for key, value in section.items():
        if isinstance(value, str) and value:
            ar = translator.translate(value, src='en', dest='ar').text
            ar_section[key] = ar
        elif isinstance(value, list):
            ar_section[key] = [translator.translate(v, src='en', dest='ar').text for v in value]
        else:
            ar_section[key] = value
    return ar_section


def build_lecture_object(file_path: Path, texts_en: list[str]) -> dict:
    name = file_path.stem
    num_match = re.search(r"\d+", name)
    lecture_id = f"L{num_match.group()}" if num_match else f"L_{name}"
    cleaned = [clean_text(t) for t in texts_en]
    en_sections = group_into_sections(cleaned)
    translator = Translator()
    ar_sections = [translate_section(sec, translator) for sec in en_sections]
    lecture = {
        "id": lecture_id,
        "source_file": str(file_path),
        "title_en": name,
        "title_ar": translator.translate(name, src='en', dest='ar').text,
        "sections_en": en_sections,
        "sections_ar": ar_sections,
    }
    return lecture


def main():
    parser = argparse.ArgumentParser(description="Full pipeline for bilingual study platform")
    parser.add_argument("project_root", nargs="?", default=os.getcwd(), help="Root folder containing lecture files")
    args = parser.parse_args()
    root = Path(args.project_root)
    print(f"Scanning {root} for lecture files …")
    files = discover_files(root)
    print(f"Found {len(files)} files.")
    lectures = []
    for f in files:
        ext = f.suffix.lower()
        if ext == ".pptx":
            texts = extract_from_pptx(f)
        elif ext == ".pdf":
            texts = extract_from_pdf(f)
        elif ext == ".docx":
            texts = extract_from_docx(f)
        else:
            continue
        if not texts:
            continue
        lecture_obj = build_lecture_object(f, texts)
        lectures.append(lecture_obj)
        print(f"Processed {f.name} → {lecture_obj['id']}, {len(texts)} raw slides/paragraphs.")

    def sort_key(l):
        m = re.search(r"\d+", l["id"])
        return int(m.group()) if m else 0
    lectures.sort(key=sort_key)

    content = {"lectures": lectures}
    output_path = root / "content.json"
    with open(output_path, "w", encoding="utf-8") as out:
        json.dump(content, out, ensure_ascii=False, indent=2)
    print(f"✅ Content written to {output_path}")

    generate_script = root / "generate_site.py"
    if generate_script.exists():
        print("Running generate_site.py …")
        subprocess.check_call([sys.executable, str(generate_script), str(output_path)])
    else:
        print("⚠️ generate_site.py not found – please add it to the project folder.")

if __name__ == "__main__":
    main()
